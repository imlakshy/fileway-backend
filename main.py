import traceback
from fastapi import FastAPI, HTTPException, Request, UploadFile, File, Form
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import fitz  # PyMuPDF
from PIL import Image, ImageOps
import requests
import io
import zipfile
from fastapi.responses import HTMLResponse

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

HTML = """
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>FileWay Backend</title>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;700&display=swap" rel="stylesheet">
  <style>
    body {
      margin: 0;
      padding: 0;
      min-height: 100vh;
      display: flex;
      align-items: center;
      justify-content: center;
      background: #121212;
      font-family: 'Inter', Arial, sans-serif;
    }
    .container {
      background: #000;
      border-radius: 18px;
      box-shadow: 0 8px 32px rgba(20, 20, 40, 0.32);
      padding: 48px 36px;
      max-width: 420px;
      text-align: center;
    }
    h1 {
      font-size: 2.2rem;
      font-weight: 700;
      margin-bottom: 18px;
      color: #74a1af;
      letter-spacing: -1px;
      line-height: 1.1;
    }
    .fileway {
      font-size: 3.6rem;
      font-weight: 800;
      color: #943900;
      letter-spacing: -2px;
      display: inline-block;
      margin: 10px 0 4px 0;
    }
    .backend {
      color: #b8b8b8;
    }
    p {
      font-size: 1.1rem;
      color: #b8b8b8;
      margin: 16px 0;
    }
    a {
      color: #ff9800;
      text-decoration: none;
      font-weight: 600;
      transition: color 0.2s;
    }
    a:hover {
      color: #7a3306;
      text-decoration: underline;
    }
    .status {
      display: inline-block;
      background: #1a1a1a;
      color: #943900;
      font-weight: 600;
      border-radius: 8px;
      padding: 4px 12px;
      margin-bottom: 8px;
      font-size: 1rem;
      box-shadow: 0 2px 8px rgba(122,51,6,0.10);
    }
    .tagline {
      font-size: 1.08rem;
      color: #ffb26b;
      font-weight: 500;
      margin: 18px 0 0 0;
      letter-spacing: 0.01em;
      background: #1a1a1a;
      border-radius: 8px;
      padding: 10px 0 10px 0;
      box-shadow: 0 2px 8px rgba(122,51,6,0.10);
      text-align: center;
      transition: background 0.2s;
    }
    .tagline .emoji {
      font-size: 1.2em;
      margin-left: 4px;
    }
  </style>
</head>
<body>
  <div class="container">
    <h1>
      Welcome to <br>
      <span class="fileway">FileWay</span><br>
      backend
    </h1>
    <div class="status">Backend is running ✅</div>
    <p>Go to <a href="https://fileway.vercel.app/" target="_blank">Frontend</a></p>
    <div class="tagline">One-stop solution for all your file types <span class="emoji">⚡🛠️✨</span></div>
  </div>
</body>
</html>
"""




@app.get("/", response_class=HTMLResponse)
async def root():
    return HTML

import time

START_TIME = time.time()

@app.get("/health")
def health():
    uptime = time.time() - START_TIME
    return {
        "status": "ok",
        "uptime": uptime
    }


@app.post("/merge-pdfs")
async def merge_pdfs(request: Request):
    try:
        body = await request.json()
        pdf_urls = body.get("pdf_urls", [])

        if not pdf_urls:
            raise HTTPException(status_code=400, detail="No PDF URLs provided.")

        merger = PdfMerger()

        for url in pdf_urls:
            print(f"Downloading: {url}")
            response = requests.get(url)
            if response.status_code == 200:
                pdf_stream = io.BytesIO(response.content)
                merger.append(pdf_stream)
            else:
                print(f"Failed to download: {url}")

        output_pdf = io.BytesIO()
        merger.write(output_pdf)
        merger.close()
        output_pdf.seek(0)

        return StreamingResponse(output_pdf, media_type="application/pdf", headers={
            "Content-Disposition": "inline; filename=merged.pdf"
        })

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error merging PDFs: {str(e)}")
    
@app.post("/unlock-pdf")
async def unlock_pdf(request: Request):
    try:
        body = await request.json()
        pdf_url = body.get("pdf_urls")
        password = body.get("password")

        if not pdf_url:
            raise HTTPException(status_code=400, detail="URL are required.")
        if  not password:
            raise HTTPException(status_code=400, detail="password are required.")
        
        if isinstance(pdf_url, list):
         pdf_url = pdf_url[0]

        print(f"Downloading: {pdf_url}")
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")

        pdf_stream = io.BytesIO(response.content)
        reader = PdfReader(pdf_stream)

        if not reader.is_encrypted:
            raise HTTPException(status_code=400, detail="PDF is not encrypted.")

        if reader.decrypt(password) == 0:
            raise HTTPException(status_code=401, detail="Incorrect password.")

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output_pdf = io.BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)

        return StreamingResponse(output_pdf, media_type="application/pdf", headers={
            "Content-Disposition": "inline; filename=unlocked.pdf"
        })

    except Exception as e:
        traceback.print_exc()  # 👈 This will print the full error in the terminal
        raise HTTPException(status_code=500, detail=f"Error unlocking PDF: {str(e)}")
    
@app.post("/split-pdf")
async def split_pdf(request: Request):
    try:
        body = await request.json()
        pdf_url = body.get("pdf_urls") # type: ignore
        start = int(body.get("start"))
        end = int(body.get("end"))

        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing url")
        if start is None or end is None:
            raise HTTPException(status_code=400, detail="Missing page range.")
        if isinstance(pdf_url, list):
         pdf_url = pdf_url[0]  #
        
        

        # Download the PDF
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")

        pdf_bytes = io.BytesIO(response.content)
        reader = PdfReader(pdf_bytes)

        if start < 1 or end > len(reader.pages) or start > end:
            raise HTTPException(status_code=400, detail="Invalid page range.")

        writer = PdfWriter()
        for i in range(start - 1, end):  # PDF page numbers are 0-based
            writer.add_page(reader.pages[i])

        output_pdf = io.BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)

        return StreamingResponse(output_pdf, media_type="application/pdf", headers={
            "Content-Disposition": "inline; filename=split_pages.pdf"
        })

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error splitting PDF: {str(e)}")
    
@app.post("/dark-mode-pdf")
async def convert_to_dark_mode(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL")
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]  # ✅ Take the first item

        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF")

        input_pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=input_pdf_stream, filetype="pdf")

        dark_pdf = fitz.open()
        for page in doc:
            pix = page.get_pixmap(dpi=150) # type: ignore
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples) # type: ignore

            dark_img = ImageOps.invert(img)

            img_buffer = io.BytesIO()
            dark_img.save(img_buffer, format="PNG")
            img_buffer.seek(0)

            rect = fitz.Rect(0, 0, pix.width, pix.height)
            pdf_page = dark_pdf.new_page(width=rect.width, height=rect.height) # type: ignore
            pdf_page.insert_image(rect, stream=img_buffer.read())

        output_buffer = io.BytesIO()
        dark_pdf.save(output_buffer)
        dark_pdf.close()
        output_buffer.seek(0)

        return StreamingResponse(output_buffer, media_type="application/pdf", headers={
            "Content-Disposition": "inline; filename=dark_mode.pdf"
        })

    except Exception as e:
        print(f"[ERROR] {str(e)}")  # 🧠 Print error in console
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")
    
@app.post("/compress-pdf")
async def compress_pdf(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        target_kb = float(data.get("target_kb", 500))

        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]

        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")

        original_pdf = fitz.open(stream=response.content, filetype="pdf")

        best_pdf_bytes = None
        best_quality = None
        best_size_kb = None

        for quality in range(80, 5, -5):  # Try different compression qualities
            compressed_pdf = fitz.open()
            for page in original_pdf:
                pix = page.get_pixmap(dpi=100)  # type: ignore # Lower DPI = smaller file
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples) # type: ignore
                img_io = io.BytesIO()
                img.save(img_io, format="JPEG", quality=quality)
                img_io.seek(0)

                rect = fitz.Rect(0, 0, pix.width, pix.height)
                new_page = compressed_pdf.new_page(width=rect.width, height=rect.height) # type: ignore
                new_page.insert_image(rect, stream=img_io.read())

            output_stream = io.BytesIO()
            compressed_pdf.save(output_stream)
            compressed_pdf.close()

            size_kb = len(output_stream.getvalue()) / 1024

            print(f"📦 Quality {quality}: {int(size_kb)} KB")

            # Store the best compressed result
            if best_size_kb is None or size_kb < best_size_kb:
                best_size_kb = size_kb
                best_pdf_bytes = output_stream.getvalue()
                best_quality = quality

            if size_kb <= target_kb:
                output_stream.seek(0)
                return StreamingResponse(io.BytesIO(best_pdf_bytes), media_type="application/pdf", headers={ # type: ignore
                    "Content-Disposition": f"inline; filename=compressed_q{quality}.pdf"
                })

        # If target was not met, return info about min possible
        return JSONResponse({
            "message": "❌ Cannot compress to desired size.",
            "min_possible_kb": round(best_size_kb, 2), # type: ignore
            "best_quality": best_quality
        }, status_code=200)

    except Exception as e:
        print("❌ Error:", e)
        raise HTTPException(status_code=500, detail=str(e))
    
@app.post("/encrypt-pdf")
async def encrypt_pdf(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        password = data.get("password")

        if not pdf_url or not password:
            raise HTTPException(status_code=400, detail="Missing PDF URL or password.")

        # Handle case where URL is a list
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]

        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")

        # Read the PDF
        input_stream = io.BytesIO(response.content)
        reader = PdfReader(input_stream)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Set encryption
        writer.encrypt(user_password=password)

        output_stream = io.BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)

        return StreamingResponse(output_stream, media_type="application/pdf", headers={
            "Content-Disposition": "inline; filename=protected.pdf"
        })

    except Exception as e:
        print("❌ Error:", e)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/pdf-to-images")
async def pdf_to_images(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL.")
        
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]
        
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")
        
        pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        num_pages = len(doc)
        
        # If only 1 page, return single JPG file
        if num_pages == 1:
            page = doc[0]
            pix = page.get_pixmap(dpi=150)  # type: ignore
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  # type: ignore
            
            img_buffer = io.BytesIO()
            img.save(img_buffer, format="JPEG")
            img_buffer.seek(0)
            
            doc.close()
            
            return StreamingResponse(img_buffer, media_type="image/jpeg", headers={
                "Content-Disposition": "inline; filename=page_1.jpg"
            })
        
        # If multiple pages, return zip file
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for page_num in range(num_pages):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=150)  # type: ignore
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  # type: ignore
                
                img_buffer = io.BytesIO()
                img.save(img_buffer, format="JPEG")
                img_buffer.seek(0)
                
                zip_file.writestr(f"page_{page_num + 1}.jpg", img_buffer.getvalue())
        
        doc.close()
        zip_buffer.seek(0)
        
        return StreamingResponse(zip_buffer, media_type="application/zip", headers={
            "Content-Disposition": "attachment; filename=pdf_images.zip"
        })
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting PDF to images: {str(e)}")

@app.post("/pdf-to-word")
async def pdf_to_word(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL.")
        
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]
        
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")
        
        pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        from docx import Document
        
        word_doc = Document()
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                word_doc.add_paragraph(text)
                if page_num < len(doc) - 1:
                    word_doc.add_page_break()
        
        doc.close()
        
        output_buffer = io.BytesIO()
        word_doc.save(output_buffer)
        output_buffer.seek(0)
        
        return StreamingResponse(output_buffer, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers={
            "Content-Disposition": "attachment; filename=converted.docx"
        })
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting PDF to Word: {str(e)}")

@app.post("/pdf-to-excel")
async def pdf_to_excel(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL.")
        
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]
        
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")
        
        pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        row_num = 1
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=row_num, column=1, value=line.strip())
                        row_num += 1
                
                ws.cell(row=row_num, column=1, value=f"--- Page {page_num + 1} ---")
                row_num += 1
        
        doc.close()
        
        output_buffer = io.BytesIO()
        wb.save(output_buffer)
        output_buffer.seek(0)
        
        return StreamingResponse(output_buffer, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers={
            "Content-Disposition": "attachment; filename=converted.xlsx"
        })
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting PDF to Excel: {str(e)}")

@app.post("/pdf-to-powerpoint")
async def pdf_to_powerpoint(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL.")
        
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]
        
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")
        
        pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = f"Page {page_num + 1}"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = text[:1000] if len(text) > 1000 else text
            
            if len(text) > 1000:
                p = tf.add_paragraph()
                p.text = "...(content truncated)"
        
        doc.close()
        
        output_buffer = io.BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)
        
        return StreamingResponse(output_buffer, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={
            "Content-Disposition": "attachment; filename=converted.pptx"
        })
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting PDF to PowerPoint: {str(e)}")

@app.post("/pdf-to-text")
async def pdf_to_text(request: Request):
    try:
        data = await request.json()
        pdf_url = data.get("pdf_urls")
        
        if not pdf_url:
            raise HTTPException(status_code=400, detail="Missing PDF URL.")
        
        if isinstance(pdf_url, list):
            pdf_url = pdf_url[0]
        
        response = requests.get(pdf_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download PDF.")
        
        pdf_stream = io.BytesIO(response.content)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        text_content = ""
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            text_content += f"--- Page {page_num + 1} ---\n\n"
            text_content += text + "\n\n"
        
        doc.close()
        
        text_buffer = io.BytesIO(text_content.encode('utf-8'))
        text_buffer.seek(0)
        
        return StreamingResponse(text_buffer, media_type="text/plain", headers={
            "Content-Disposition": "attachment; filename=converted.txt"
        })
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting PDF to text: {str(e)}")

###==================================================================###

@app.post("/changeImgExt")
async def convert_image_urls(request: Request):
    try:
        data = await request.json()
        img_urls = data.get("img_urls", [])
        desired_ext = data.get("UserDesiredConvertedExtension")

        if not img_urls or not isinstance(img_urls, list):
            raise HTTPException(status_code=400, detail="img_urls must be a non-empty array.")

        # Accept formats like: ["JPEG","PNG","WEBP","PDF","GIF","BMP","TIFF","ICO","PPM","EPS"]
        desired = (desired_ext or "").strip().upper().lstrip(".")
        allowed_formats = {
            "JPEG": {"ext": "jpg", "mime": "image/jpeg"},
            "JPG": {"ext": "jpg", "mime": "image/jpeg"},
            "PNG": {"ext": "png", "mime": "image/png"},
            "WEBP": {"ext": "webp", "mime": "image/webp"},
            "PDF": {"ext": "pdf", "mime": "application/pdf"},
            "GIF": {"ext": "gif", "mime": "image/gif"},
            "BMP": {"ext": "bmp", "mime": "image/bmp"},
            "TIFF": {"ext": "tiff", "mime": "image/tiff"},
            "TIF": {"ext": "tiff", "mime": "image/tiff"},
            "ICO": {"ext": "ico", "mime": "image/x-icon"},
            "PPM": {"ext": "ppm", "mime": "image/x-portable-pixmap"},
            "EPS": {"ext": "eps", "mime": "application/postscript"},
        }

        if desired not in allowed_formats:
            raise HTTPException(
                status_code=400,
                detail="Unsupported format. Use one of: JPEG, PNG, WEBP, PDF, GIF, BMP, TIFF, ICO, PPM, EPS",
            )

        target_format = "JPEG" if desired == "JPG" else desired
        out_ext = allowed_formats[desired]["ext"]
        out_mime = allowed_formats[desired]["mime"]

        def _convert_one(file_bytes: bytes) -> bytes:
            with Image.open(io.BytesIO(file_bytes)) as img:
                if target_format == "JPEG" and img.mode in ("RGBA", "LA", "P"):
                    img = img.convert("RGB")
                out = io.BytesIO()
                img.save(out, format=target_format)
                return out.getvalue()

        def _download(url: str) -> bytes:
            resp = requests.get(url)
            if resp.status_code != 200:
                raise HTTPException(status_code=400, detail=f"Failed to download image: {url}")
            return resp.content

        # Single URL -> return image directly
        if len(img_urls) == 1:
            img_bytes = _download(img_urls[0])
            converted = _convert_one(img_bytes)
            return StreamingResponse(
                io.BytesIO(converted),
                media_type=out_mime,
                headers={"Content-Disposition": f"attachment; filename=converted.{out_ext}"},
            )

        # Multiple URLs -> zip
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for idx, url in enumerate(img_urls, start=1):
                img_bytes = _download(url)
                converted = _convert_one(img_bytes)
                zf.writestr(f"image_{idx}.{out_ext}", converted)

        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=converted_images.zip"},
        )

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting image URLs: {str(e)}")


@app.post("/resizeImgByKB")
async def resize_img_by_kb(request: Request):
    """
    Accepts JSON:
      {
        "img_urls": ["https://.../a.png", "https://.../b.webp"],
        "sizeInKB": 200
      }
    Output:
      - Always converts to JPEG/JPG
      - If 1 url: returns a single JPG
      - If >1 urls: returns a zip of JPGs
      - Each output file will be <= sizeInKB (best-effort; will error if impossible)
    """
    try:
        data = await request.json()
        img_urls = data.get("img_urls", [])
        size_kb = data.get("sizeInKB")

        if not img_urls or not isinstance(img_urls, list):
            raise HTTPException(status_code=400, detail="img_urls must be a non-empty array.")

        if size_kb is None:
            raise HTTPException(status_code=400, detail="sizeInKB is required.")

        try:
            target_kb = float(size_kb)
        except Exception:
            raise HTTPException(status_code=400, detail="sizeInKB must be a number.")

        if target_kb <= 0:
            raise HTTPException(status_code=400, detail="sizeInKB must be > 0.")

        target_bytes = int(target_kb * 1024)

        def _download(url: str) -> bytes:
            resp = requests.get(url)
            if resp.status_code != 200:
                raise HTTPException(status_code=400, detail=f"Failed to download image: {url}")
            return resp.content

        def _save_jpeg(
            img: Image.Image,
            quality: int,
            *,
            optimize: bool,
            subsampling: int,
        ) -> bytes:
            out = io.BytesIO()
            img.save(
                out,
                format="JPEG",
                quality=quality,
                optimize=optimize,
                subsampling=subsampling,
            )
            return out.getvalue()

        def _compress_to_target(file_bytes: bytes) -> bytes:
            with Image.open(io.BytesIO(file_bytes)) as im:
                # Normalize to RGB for JPEG
                img = im.convert("RGB")
                tolerance = 0.05  # aim within ±5%
                lower = int(target_bytes * (1 - tolerance))
                upper = int(target_bytes * (1 + tolerance))

                def encode_at(img_: Image.Image, q: int) -> bytes:
                    # For smaller outputs, optimize helps; for larger outputs, disable optimize.
                    want_bigger = target_bytes > 0
                    optimize = not want_bigger
                    # subsampling=2 is smaller, subsampling=0 is larger/better quality
                    subsampling = 2 if optimize else 0
                    return _save_jpeg(img_, q, optimize=optimize, subsampling=subsampling)

                # Helper: find closest by varying quality for a given image size
                def best_by_quality(img_: Image.Image, q_min: int, q_max: int) -> bytes:
                    best_bytes = None
                    best_diff = None
                    for q in range(q_max, q_min - 1, -1):
                        b = _save_jpeg(
                            img_,
                            q,
                            optimize=False if target_bytes >= lower else True,
                            subsampling=0 if target_bytes >= lower else 2,
                        )
                        d = abs(len(b) - target_bytes)
                        if best_diff is None or d < best_diff:
                            best_bytes, best_diff = b, d
                        if lower <= len(b) <= upper:
                            return b
                    return best_bytes  # type: ignore

                # First try at original size: sweep quality to get close
                current_best = best_by_quality(img, 10, 95)
                if lower <= len(current_best) <= upper:
                    return current_best

                # If we need to SHRINK (target smaller than current best), downscale gradually
                if len(current_best) > upper:
                    w, h = img.size
                    scale = 0.9
                    while w > 50 and h > 50:
                        nw, nh = max(1, int(w * scale)), max(1, int(h * scale))
                        resized = img.resize((nw, nh), Image.LANCZOS)
                        candidate = best_by_quality(resized, 10, 95)
                        if lower <= len(candidate) <= upper:
                            return candidate
                        current_best = candidate
                        if len(candidate) <= upper:
                            # best effort: closest under target
                            return candidate
                        w, h = nw, nh
                    raise HTTPException(
                        status_code=400,
                        detail="Cannot compress image to the requested sizeInKB.",
                    )

                # If we need to ENLARGE (target larger than current best), upscale gradually
                w, h = img.size
                scale = 1.15
                max_scale = 4.0
                total_scale = 1.0
                while total_scale < max_scale:
                    if len(current_best) >= lower:
                        break
                    total_scale *= scale
                    nw, nh = max(1, int(w * total_scale)), max(1, int(h * total_scale))
                    up = img.resize((nw, nh), Image.LANCZOS)
                    # Use settings that generally increase JPEG size
                    candidate = _save_jpeg(up, 95, optimize=False, subsampling=0)
                    current_best = candidate
                    if lower <= len(candidate) <= upper:
                        return candidate
                    if len(candidate) >= lower:
                        # now bracketed-ish; refine with quality sweep to get closer
                        return best_by_quality(up, 10, 95)

                # If still smaller than target even after upscaling, return best effort (largest we got)
                return current_best

        # Single URL -> return image directly
        if len(img_urls) == 1:
            img_bytes = _download(img_urls[0])
            converted = _compress_to_target(img_bytes)
            return StreamingResponse(
                io.BytesIO(converted),
                media_type="image/jpeg",
                headers={"Content-Disposition": "attachment; filename=compressed.jpg"},
            )

        # Multiple URLs -> zip
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for idx, url in enumerate(img_urls, start=1):
                img_bytes = _download(url)
                converted = _compress_to_target(img_bytes)
                zf.writestr(f"image_{idx}.jpg", converted)

        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=compressed_images.zip"},
        )

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error resizing images: {str(e)}")


@app.post("/resizeImgByHW")
async def resize_img_by_height_width(request: Request):
    try:
        data = await request.json()
        img_urls = data.get("img_urls", [])
        width = data.get("width")
        height = data.get("height")

        if not img_urls or not isinstance(img_urls, list):
            raise HTTPException(status_code=400, detail="img_urls must be a non-empty array.")

        if width is None or height is None:
            raise HTTPException(status_code=400, detail="width and height are required.")

        try:
            w = int(width)
            h = int(height)
        except Exception:
            raise HTTPException(status_code=400, detail="width and height must be numbers.")

        if w <= 0 or h <= 0:
            raise HTTPException(status_code=400, detail="width and height must be > 0.")

        def _download(url: str) -> bytes:
            resp = requests.get(url)
            if resp.status_code != 200:
                raise HTTPException(status_code=400, detail=f"Failed to download image: {url}")
            return resp.content

        # Map PIL format -> (extension, mime)
        fmt_map = {
            "JPEG": ("jpg", "image/jpeg"),
            "JPG": ("jpg", "image/jpeg"),
            "PNG": ("png", "image/png"),
            "WEBP": ("webp", "image/webp"),
            "GIF": ("gif", "image/gif"),
            "BMP": ("bmp", "image/bmp"),
            "TIFF": ("tiff", "image/tiff"),
            "ICO": ("ico", "image/x-icon"),
            "PPM": ("ppm", "image/x-portable-pixmap"),
            "EPS": ("eps", "application/postscript"),
            "PDF": ("pdf", "application/pdf"),
        }

        def _resize_keep_format(file_bytes: bytes) -> tuple[bytes, str, str]:
            with Image.open(io.BytesIO(file_bytes)) as im:
                pil_format = (im.format or "PNG").upper()
                out_ext, out_mime = fmt_map.get(pil_format, ("png", "image/png"))

                resized = im.resize((w, h), Image.LANCZOS)

                # Handle modes when saving to formats that don't support alpha
                save_format = pil_format if pil_format in fmt_map else "PNG"
                if save_format in ("JPEG", "JPG") and resized.mode in ("RGBA", "LA", "P"):
                    resized = resized.convert("RGB")

                out = io.BytesIO()
                resized.save(out, format=save_format)
                return out.getvalue(), out_ext, out_mime

        # Single URL -> return file directly
        if len(img_urls) == 1:
            img_bytes = _download(img_urls[0])
            resized_bytes, out_ext, out_mime = _resize_keep_format(img_bytes)
            return StreamingResponse(
                io.BytesIO(resized_bytes),
                media_type=out_mime,
                headers={"Content-Disposition": f"attachment; filename=resized.{out_ext}"},
            )

        # Multiple URLs -> zip
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for idx, url in enumerate(img_urls, start=1):
                img_bytes = _download(url)
                resized_bytes, out_ext, _ = _resize_keep_format(img_bytes)
                zf.writestr(f"image_{idx}.{out_ext}", resized_bytes)

        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=resized_images.zip"},
        )

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error resizing images by height/width: {str(e)}")
